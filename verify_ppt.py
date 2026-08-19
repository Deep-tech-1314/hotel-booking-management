from pptx import Presentation
import re

pptx_path = r"C:\Users\LENOVO\Downloads\Major Project-I (01CE0716) - PPT format.pptx"
output_path = r"d:\Hotel Booking Management\ppt_modified_structure.txt"

try:
    print(f"Loading modified PPTX from: {pptx_path}")
    prs = Presentation(pptx_path)
    print(f"Loaded successfully! Slide count: {len(prs.slides)}")
    
    with open(output_path, 'w', encoding='utf-8') as out:
        out.write(f"SLIDE_COUNT: {len(prs.slides)}\n\n")
        
        for idx, slide in enumerate(prs.slides):
            out.write(f"=== Slide {idx + 1} (Layout: {slide.slide_layout.name}) ===\n")
            # Extract title if present
            title_text = ""
            for shape in slide.shapes:
                if shape.is_placeholder and (shape.placeholder_format.type == 1 or shape.placeholder_format.type == 3):
                    if shape.has_text_frame:
                        title_text = shape.text_frame.text
                        break
            
            if title_text:
                out.write(f"Title: {title_text}\n")
                
            # Extract all text from shapes
            all_texts = []
            for shape in slide.shapes:
                if shape.has_text_frame:
                    text = shape.text_frame.text.strip()
                    if text and text != title_text:
                        all_texts.append(text)
            
            if all_texts:
                out.write("Texts:\n" + "\n---\n".join(all_texts) + "\n")
            out.write("\n" + "="*40 + "\n\n")
            
    print(f"Verification results written to {output_path}")
except Exception as e:
    print(f"Verification failed: {e}")
    with open(output_path, 'w', encoding='utf-8') as out:
        out.write(f"ERROR: {e}\n")
