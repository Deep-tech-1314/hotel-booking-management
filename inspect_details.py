import sys
import subprocess

try:
    from pptx import Presentation
except ImportError:
    print("Installing python-pptx library...")
    subprocess.check_call([sys.executable, "-m", "pip", "install", "python-pptx"])
    from pptx import Presentation

pptx_path = r"C:\Users\LENOVO\Downloads\Major Project-I (01CE0716) - PPT format.pptx"
output_path = r"d:\Hotel Booking Management\ppt_detailed_structure.txt"

try:
    prs = Presentation(pptx_path)
    with open(output_path, 'w', encoding='utf-8') as out:
        out.write(f"Slide Width: {prs.slide_width}\n")
        out.write(f"Slide Height: {prs.slide_height}\n")
        out.write(f"Number of Slides: {len(prs.slides)}\n\n")
        
        # Log layouts
        out.write("=== Slide Layouts in Template ===\n")
        for idx, layout in enumerate(prs.slide_layouts):
            out.write(f"Layout Index {idx}: {layout.name}\n")
            for shape in layout.shapes:
                if shape.is_placeholder:
                    out.write(f"  Placeholder: name='{shape.name}', type={shape.placeholder_format.type}, idx={shape.placeholder_format.idx}\n")
        out.write("\n")
        
        # Log each slide's shapes
        for s_idx, slide in enumerate(prs.slides):
            out.write(f"=== Slide {s_idx + 1} (Layout: {slide.slide_layout.name}) ===\n")
            for shape in slide.shapes:
                shape_type = shape.shape_type
                text = ""
                if shape.has_text_frame:
                    text = shape.text_frame.text
                is_placeholder = shape.is_placeholder
                ph_type = shape.placeholder_format.type if is_placeholder else "N/A"
                ph_idx = shape.placeholder_format.idx if is_placeholder else "N/A"
                out.write(f"Shape ID {shape.shape_id}: '{shape.name}' | Type: {shape_type} | Placeholder: {is_placeholder} (type={ph_type}, idx={ph_idx})\n")
                if text:
                    out.write(f"  Text: {text}\n")
            out.write("\n")
            
    print(f"Detailed structure written to {output_path}")
except Exception as e:
    with open(output_path, 'w', encoding='utf-8') as out:
        out.write(f"ERROR: {e}\n")
    print(f"Error: {e}")
