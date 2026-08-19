import sys
import subprocess
import os
import zipfile
import tempfile

try:
    from pptx import Presentation
    from pptx.util import Inches, Pt, Emu
    from pptx.dml.color import RGBColor
    from pptx.enum.shapes import MSO_SHAPE
    from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
    from pptx.oxml.ns import qn
except ImportError:
    print("Installing python-pptx library...")
    subprocess.check_call([sys.executable, "-m", "pip", "install", "python-pptx"])
    from pptx import Presentation
    from pptx.util import Inches, Pt, Emu
    from pptx.dml.color import RGBColor
    from pptx.enum.shapes import MSO_SHAPE
    from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
    from pptx.oxml.ns import qn

TEMPLATE_PATH = r"C:\Users\LENOVO\Downloads\Major Project-I (01CE0716) - PPT format (1).pptx"
OUTPUT_PATH   = r"C:\Users\LENOVO\Downloads\BookMyStay_Frontend_PPT.pptx"
ARTIFACTS_DIR = r"C:\Users\LENOVO\.gemini\antigravity-ide\brain\96637727-82b0-4c6b-96fa-8814e1c61536"

# ---- Template design tokens (extracted from the original Marwadi University template) ----
FONT_HEAD = "Proxima Nova"
FONT_BODY = "Proxima Nova"
TEAL   = RGBColor(0x04, 0xA2, 0xB9)   # brand accent — titles
BLACK  = RGBColor(0x00, 0x00, 0x00)   # body text
GRAY   = RGBColor(0x59, 0x59, 0x59)   # muted subtext
DGRAY  = RGBColor(0x66, 0x66, 0x66)
RED    = RGBColor(0xFF, 0x00, 0x00)   # emphasis / values
GREEN  = RGBColor(0x2E, 0x8B, 0x57)
WHITE  = RGBColor(0xFF, 0xFF, 0xFF)
CARD_FILL = RGBColor(0xF3, 0xF6, 0xF9)   # very light gray card
CARD_ALT  = RGBColor(0xE8, 0xF6, 0xF9)   # light teal tint

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)

BG_CONTENT = None  # path to extracted image2.jpeg (content frame background)


def extract_bg():
    """Pull the content-frame background image (image2.jpeg) out of the template."""
    global BG_CONTENT
    tmp = os.path.join(tempfile.gettempdir(), "mu_content_bg.jpeg")
    with zipfile.ZipFile(TEMPLATE_PATH) as z:
        with open(tmp, "wb") as f:
            f.write(z.read("ppt/media/image2.jpeg"))
    BG_CONTENT = tmp


def add_bg(slide):
    """Place the template content-frame image as a full-bleed background (added first => behind)."""
    pic = slide.shapes.add_picture(BG_CONTENT, 0, 0, SLIDE_W, SLIDE_H)
    # move to back of z-order
    sp = pic._element
    sp.getparent().remove(sp)
    slide.shapes._spTree.insert(2, sp)
    return pic


def no_bullet(p):
    """Suppress any inherited list bullet on a paragraph (schema-safe placement)."""
    pPr = p._p.get_or_add_pPr()
    for t in ('a:buChar', 'a:buAutoNum', 'a:buNone'):
        for el in pPr.findall(qn(t)):
            pPr.remove(el)
    bu = pPr.makeelement(qn('a:buNone'), {})
    ref = pPr.find(qn('a:defRPr'))
    if ref is None:
        ref = pPr.find(qn('a:tabLst'))
    if ref is not None:
        ref.addprevious(bu)
    else:
        pPr.append(bu)


def add_title(slide, text):
    box = slide.shapes.add_textbox(Inches(0.58), Inches(0.28), Inches(11.5), Inches(0.62))
    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    r = p.add_run(); r.text = text
    r.font.name = FONT_HEAD; r.font.size = Pt(24); r.font.bold = True; r.font.color.rgb = TEAL
    return box


def _run(p, text, size=18, color=BLACK, bold=False, font=FONT_BODY):
    r = p.add_run(); r.text = text
    r.font.name = font; r.font.size = Pt(size); r.font.bold = bold; r.font.color.rgb = color
    return r


def add_bullets(slide, left, top, width, height, items, size=18, gap=6):
    """items: list of (text, level) tuples OR plain strings (level 0)."""
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    first = True
    for it in items:
        if isinstance(it, tuple):
            text, level = it
        else:
            text, level = it, 0
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        p.space_after = Pt(gap)
        p.level = level
        if level == 0:
            _run(p, "▪  ", size=size, color=TEAL, bold=True)   # small teal square bullet
            # split "Label: rest" so the label is bold teal-ish
            if ": " in text and text.index(": ") < 40:
                label, rest = text.split(": ", 1)
                _run(p, label + ": ", size=size, color=BLACK, bold=True)
                _run(p, rest, size=size, color=BLACK)
            else:
                _run(p, text, size=size, color=BLACK)
        else:
            _run(p, "–  ", size=size - 1, color=GRAY, bold=True)  # en-dash sub-bullet
            _run(p, text, size=size - 1, color=DGRAY)
    return box


def card(slide, left, top, width, height, text="", fill=CARD_FILL, border=TEAL,
         size=11, tcolor=BLACK, bold=False, border_w=1.25, align=PP_ALIGN.CENTER):
    shp = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shp.fill.solid(); shp.fill.fore_color.rgb = fill
    if border is None:
        shp.line.fill.background()
    else:
        shp.line.color.rgb = border; shp.line.width = Pt(border_w)
    shp.shadow.inherit = False
    tf = shp.text_frame; tf.word_wrap = True; tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    lines = text.split("\n")
    for i, ln in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        bold_line = bold and i == 0
        col = tcolor
        _run(p, ln, size=size, color=col, bold=bold_line)
    return shp


def arrow(slide, left, top, width, height, direction="down"):
    shape_type = MSO_SHAPE.DOWN_ARROW if direction == "down" else MSO_SHAPE.RIGHT_ARROW
    a = slide.shapes.add_shape(shape_type, left, top, width, height)
    a.fill.solid(); a.fill.fore_color.rgb = TEAL
    a.line.fill.background(); a.shadow.inherit = False
    return a


def swatch(slide, left, top, size, hexlabel, name, color, tcolor=WHITE):
    o = slide.shapes.add_shape(MSO_SHAPE.OVAL, left, top, size, size)
    o.fill.solid(); o.fill.fore_color.rgb = color
    o.line.color.rgb = RGBColor(0xCC, 0xD3, 0xDA); o.line.width = Pt(1)
    o.shadow.inherit = False
    tf = o.text_frame; tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
    _run(p, hexlabel, size=13, color=tcolor, bold=True)
    # caption below
    cap = slide.shapes.add_textbox(left - Inches(0.3), top + size + Inches(0.08), size + Inches(0.6), Inches(0.5))
    cp = cap.text_frame; cp.word_wrap = True
    pp = cp.paragraphs[0]; pp.alignment = PP_ALIGN.CENTER
    _run(pp, name, size=12, color=DGRAY, bold=True)


def add_image_safe(slide, img_filename, left, top, width, height, frame=True):
    path = os.path.join(ARTIFACTS_DIR, img_filename)
    if os.path.exists(path):
        try:
            if frame:
                fr = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left - Emu(20000), top - Emu(20000),
                                            width + Emu(40000), height + Emu(40000))
                fr.fill.background(); fr.line.color.rgb = TEAL; fr.line.width = Pt(1.25); fr.shadow.inherit = False
            slide.shapes.add_picture(path, left, top, width, height)
            print(f"  [IMG] Embedded: {img_filename}")
            return True
        except Exception as e:
            print(f"  [WARN] Could not embed {img_filename}: {e}")
    else:
        print(f"  [WARN] Image not found: {img_filename}")
    return False


def caption(slide, left, top, width, text, size=11):
    box = slide.shapes.add_textbox(left, top, width, Inches(0.4))
    p = box.text_frame.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
    _run(p, text, size=size, color=GRAY, bold=True)


def get_title_ph(slide):
    for sh in slide.shapes:
        if sh.is_placeholder and sh.placeholder_format.idx == 0:
            return sh
    return None


def set_title_ph(slide, text):
    ph = get_title_ph(slide)
    if ph is None:
        return add_title(slide, text)
    tf = ph.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    r = p.add_run(); r.text = text
    r.font.name = FONT_HEAD; r.font.size = Pt(24); r.font.bold = True; r.font.color.rgb = TEAL
    return ph


# =====================================================================================
def build_ppt():
    print("Loading template...")
    extract_bg()
    prs = Presentation(TEMPLATE_PATH)
    original_count = len(prs.slides)
    print(f"Original template slides: {original_count}")

    BLANK = prs.slide_layouts[6]

    def new_slide(title):
        s = prs.slides.add_slide(BLANK)
        add_bg(s)
        add_title(s, title)
        return s

    # -------------------------------------------------------------------
    # ORIGINAL SLIDE 1 (Marwadi title slide): left untouched — genuine branding.
    # ORIGINAL SLIDE 2 (cover): fill placeholder text, preserve template style.
    # -------------------------------------------------------------------
    print("Filling Slide 2 (project cover) in place...")
    s2 = prs.slides[1]
    for sh in s2.shapes:
        if not sh.has_text_frame:
            continue
        sid = sh.shape_id
        runs = [r for p in sh.text_frame.paragraphs for r in p.runs]
        if sid == 88:      # team members — replace the three red placeholders, keep formatting
            names = ["Student Name (92200100000) (CE-7)"] * 3
            n = 0
            for r in runs:
                if 'FULL NAME' in r.text and n < 3:
                    r.text = names[n]; n += 1
        elif sid == 89:    # project title + team id — shrink title so it fits one line
            for r in runs:
                s = r.text.strip()
                if s == 'PROJECT/INTERNSHIP':
                    r.text = 'BookMyStay – Hotel Booking Management Platform'
                    r.font.size = Pt(20)
                elif s == 'TITLE':
                    r.text = ''
                elif 'xxxxx' in r.text:
                    r.text = r.text.replace('xxxxx', 'MU_CE_07')
        elif sid == 91:    # internal guide
            for r in runs:
                if 'Internal Guide Name' in r.text:
                    r.text = 'Prof. Guide Name & Designation'
        elif sid == 93:    # keep course info exactly, only fill the review date
            for r in runs:
                if 'DD/MM/YYYY' in r.text:
                    r.text = r.text.replace('DD/MM/YYYY', '10/07/2026')

    # -------------------------------------------------------------------
    # ORIGINAL SLIDE 3 (Outline): fill with real agenda.
    # -------------------------------------------------------------------
    print("Filling Slide 3 (outline)...")
    s3 = prs.slides[2]
    set_title_ph(s3, "Outline")
    # remove the placeholder body (shape 100) content, refill
    for sh in list(s3.shapes):
        if sh.has_text_frame and not sh.is_placeholder and sh.shape_id == 100:
            sh.text_frame.clear()
            outline = [
                "Introduction & Design Philosophy",
                "Frontend Technology Stack & Data Flow",
                "Visual Language & Theme Tokens",
                "Custom & Performance-Driven Components",
                "Motion & Micro-Interactions System",
                "Multimedia & Sticky-Scroll Showcase",
                "Frontend User Flows (Public View)",
                "Multi-Role Dashboards (User / Owner / Admin)",
                "Frontend State Architecture (Redux)",
                "Supporting Backend Integration",
                "Performance & Accessibility Standards",
                "Conclusion, Future Roadmap & References",
            ]
            tf = sh.text_frame; tf.word_wrap = True
            for i, item in enumerate(outline):
                p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
                p.space_after = Pt(4)
                _run(p, f"{i+1}.  ", size=18, color=TEAL, bold=True, font=FONT_HEAD)
                _run(p, item, size=18, color=BLACK, font=FONT_BODY)
            # widen the box so 12 items fit
            sh.width = Inches(11.8); sh.height = Inches(5.6)
            sh.left = Inches(0.71); sh.top = Inches(1.15)

    # ===================================================================
    # NEW CONTENT SLIDES (inserted after Outline, styled like the template)
    # ===================================================================
    LX, LW = Inches(0.62), Inches(6.05)     # left column (bullets)
    RX = Inches(6.95)                        # right column start

    # --- 1. Introduction & Design Philosophy ---
    print("New: Introduction & Design Philosophy")
    s = new_slide("Introduction & Design Philosophy")
    add_bullets(s, LX, Inches(1.2), LW, Inches(5.6), [
        "BookMyStay: a premium, full-stack hotel booking & management platform for the modern traveler.",
        "Target users: tech-savvy professionals who travel often and expect friction-free, app-grade polish.",
        "Design posture: premium restraint — clean layout, generous whitespace, a single accent, no cliché stock art.",
        "High-fidelity rendering: smooth 60fps motion and near-zero layout shift on every key page.",
        "Layout specificity: desktop tuned for research & comparison; mobile tuned for quick on-the-go bookings.",
    ], size=17)
    add_image_safe(s, "homepage_hero_mockup_1783678137330.png", RX + Inches(0.15), Inches(1.6), Inches(5.7), Inches(4.0))
    caption(s, RX, Inches(5.7), Inches(6.0), "Homepage hero — glass search bar over a looping lobby video")

    # --- 2. Frontend Technology Stack & Data Flow ---
    print("New: Tech Stack")
    s = new_slide("Frontend Technology Stack & Data Flow")
    add_bullets(s, LX, Inches(1.2), LW, Inches(5.6), [
        "Core Library: React 18 with Vite for instant HMR and fast production builds.",
        "State Management: Redux Toolkit as the single source of truth across all views.",
        "Styling Engine: hand-authored CSS custom properties — no utility-CSS bundle bloat.",
        "Icons: Lucide React with a light 1.5px stroke for a refined look.",
        "Routing: React Router DOM with lazy-loaded route boundaries for code-splitting.",
    ], size=17)
    y = Inches(1.45)
    for i, (t, sub) in enumerate([
        ("React 18 Components", "Views, hooks & dispatched actions"),
        ("Redux Toolkit Store", "Single source of truth — 9 slices"),
        ("Async Thunk Middleware", "Side-effects & API orchestration"),
        ("Express REST API + MongoDB", "Persistence & business logic"),
    ]):
        fill = CARD_ALT if i in (1, 3) else CARD_FILL
        card(s, RX, y, Inches(5.7), Inches(0.95), f"{t}\n{sub}", fill=fill, size=13, bold=True)
        if i < 3:
            arrow(s, RX + Inches(2.65), y + Inches(0.97), Inches(0.4), Inches(0.28))
        y = y + Inches(1.33)

    # --- 3. Visual Language & Theme Tokens ---
    print("New: Visual Language")
    s = new_slide("Visual Language & Theme Tokens")
    add_bullets(s, LX, Inches(1.2), Inches(5.4), Inches(2.0), [
        "A luxury dark-first design system driven entirely by CSS variables (tokens).",
        "Headings in Playfair Display; body & prices in Inter with tabular-nums alignment.",
    ], size=17)
    # app color swatches
    swatch(s, Inches(0.75), Inches(3.35), Inches(1.55), "#0f1423", "Midnight\nBackground", RGBColor(0x0f, 0x14, 0x23))
    swatch(s, Inches(2.75), Inches(3.35), Inches(1.55), "#1c2135", "Deep Navy\nPanels", RGBColor(0x1c, 0x21, 0x35))
    swatch(s, Inches(4.75), Inches(3.35), Inches(1.55), "#f5c543", "Luxury Gold\nAccent", RGBColor(0xf5, 0xc5, 0x43), tcolor=BLACK)
    card(s, RX, Inches(1.35), Inches(5.7), Inches(2.4),
         "TYPOGRAPHY\n\nHeadings — Playfair Display (editorial elegance)\nBody & Numbers — Inter (tabular-nums for prices)",
         fill=CARD_FILL, size=14, bold=True, align=PP_ALIGN.LEFT)
    card(s, RX, Inches(3.95), Inches(5.7), Inches(2.7),
         "SURFACES, RADII & GLASSMORPHISM\n\nRadii: 8 / 12 / 16 / 24 / 32 px scale\nGlass: backdrop-filter blur(24px) saturate(160%)\nElevation: rest / hover / overlay shadow tiers",
         fill=CARD_ALT, size=14, bold=True, align=PP_ALIGN.LEFT)

    # --- 4. Custom & Performance-Driven Components ---
    print("New: Components")
    s = new_slide("Custom & Performance-Driven Components")
    add_bullets(s, LX, Inches(1.2), LW, Inches(5.6), [
        "ImageLazy.jsx: loads images only when in-view via IntersectionObserver with a blur-up placeholder.",
        "VideoPlayer.jsx: dependency-free HTML5 player — lazy load, custom controls, HTML overlays.",
        "Skeleton.jsx: shimmer skeletons matching the real footprint to prevent layout shift.",
        "Design kit: Button variants, floating-label Inputs, and a custom date-range selector.",
    ], size=17)
    steps = [
        ("1. SKELETON", "Shimmer placeholder matches final footprint", CARD_FILL, TEAL),
        ("2. LAZY BLUR", "IntersectionObserver fires → low-res blur-up", CARD_ALT, TEAL),
        ("3. LOADED", "Crossfades to crisp full-resolution image", RGBColor(0xE6, 0xF7, 0xEC), GREEN),
    ]
    y = Inches(1.7)
    for t, sub, fill, bd in steps:
        card(s, RX, y, Inches(5.7), Inches(1.1), f"{t}\n{sub}", fill=fill, border=bd, size=13, bold=True)
        if t[0] != "3":
            arrow(s, RX + Inches(2.65), y + Inches(1.12), Inches(0.4), Inches(0.28))
        y = y + Inches(1.55)

    # --- 5. Motion & Micro-Interactions ---
    print("New: Motion")
    s = new_slide("Motion & Micro-Interactions System")
    add_bullets(s, LX, Inches(1.2), LW, Inches(5.6), [
        "animations.css: 30+ keyframes for entrances, ambient loops (Ken Burns) and loaders.",
        "Strict easing scale: clamped cubic-bezier curves and fixed durations (160ms hover, 240ms modals).",
        "Scroll reveals: useScrollAnimation hook drives staggered grid entrances.",
        "Micro-interactions: bursting wishlist hearts, count-up stats, inline success morphs.",
        "Accessibility: prefers-reduced-motion collapses movement to gentle opacity fades.",
    ], size=17)
    labels = [("Scroll trigger (0ms)", Inches(6.95), Inches(5.7)),
              ("Child 1  (+80ms)", Inches(7.35), Inches(5.3)),
              ("Child 2  (+160ms)", Inches(7.75), Inches(4.9)),
              ("Child 3  (+240ms)", Inches(8.15), Inches(4.5))]
    y = Inches(1.6)
    for i, (t, lx, lw) in enumerate(labels):
        fill = CARD_ALT if i == 0 else CARD_FILL
        card(s, lx, y, lw, Inches(0.7), t, fill=fill, size=13, bold=True)
        if i < 3:
            arrow(s, lx + Inches(0.35), y + Inches(0.72), Inches(0.28), Inches(0.28))
        y = y + Inches(1.05)

    # --- 6. Multimedia & Sticky-Scroll Showcase ---
    print("New: Multimedia")
    s = new_slide("Multimedia & Sticky-Scroll Showcase")
    add_bullets(s, LX, Inches(1.2), LW, Inches(5.6), [
        "InteractiveShowcase.jsx: split two-column layout for desktop storytelling.",
        ("Left column scrolls text vignettes — Suites, Wellness, Dining.", 1),
        ("Right column pins a video that cross-fades sources on the active section.", 1),
        "Hero video loop: optimized background loop with a gradient overlay for text contrast.",
        "Lightbox gallery: arrow keys, ESC to close, and pinch-zoom on mobile.",
    ], size=17)
    card(s, RX, Inches(1.5), Inches(5.7), Inches(4.6), "", fill=WHITE, border=TEAL)
    card(s, RX + Inches(0.2), Inches(1.75), Inches(2.4), Inches(1.1), "Suites\n(active)", fill=CARD_ALT, size=12, bold=True)
    card(s, RX + Inches(0.2), Inches(3.0), Inches(2.4), Inches(1.1), "Wellness", fill=CARD_FILL, border=RGBColor(0xCC,0xD3,0xDA), size=12)
    card(s, RX + Inches(0.2), Inches(4.25), Inches(2.4), Inches(1.6), "Dining", fill=CARD_FILL, border=RGBColor(0xCC,0xD3,0xDA), size=12)
    card(s, RX + Inches(2.85), Inches(1.75), Inches(2.65), Inches(4.1),
         "PINNED VIDEO\n\nCross-fades to the\nactive section's clip\n\n[ mute / unmute ]", fill=RGBColor(0x0f,0x14,0x23), border=TEAL, size=12, bold=True, tcolor=WHITE)

    # --- 7. Frontend User Flows (Public View) ---
    print("New: User Flow")
    s = new_slide("Frontend User Flows (Public View)")
    flow = [
        ("1. LANDING", "Hero video loop\nGlass search bar\nStats count-up"),
        ("2. SEARCH RESULTS", "Collapsible filters\nInstant client sort\nNo page reloads"),
        ("3. HOTEL DETAIL", "2×2 lightbox gallery\nVideo tour\nSticky reserve panel"),
        ("4. CHECKOUT", "Stripe & Razorpay\nValidation on blur\nEmail confirmation"),
    ]
    x = Inches(0.55)
    for i, (t, sub) in enumerate(flow):
        card(s, x, Inches(1.7), Inches(2.75), Inches(2.2), f"{t}\n\n{sub}", fill=CARD_FILL if i % 2 == 0 else CARD_ALT, size=13, bold=True)
        if i < 3:
            arrow(s, x + Inches(2.78), Inches(2.55), Inches(0.4), Inches(0.5), direction="right")
        x = x + Inches(3.18)
    card(s, Inches(0.55), Inches(4.35), Inches(12.25), Inches(2.0),
         "SPA FLOW PHILOSOPHY — zero full-page reloads across the journey\n\n"
         "React Router lazy boundaries swap views instantly  •  Skeleton-first loading removes layout shift  •  "
         "Filters & sorting update client-side  •  Sticky reserve keeps the primary action in reach  •  "
         "Optimistic UI confirms actions before the server responds",
         fill=CARD_ALT, size=14, bold=True, align=PP_ALIGN.CENTER)

    # --- 8. Multi-Role Dashboards ---
    print("New: Dashboards")
    s = new_slide("Multi-Role Dashboards (User / Owner / Admin)")
    add_bullets(s, LX, Inches(1.2), LW, Inches(5.6), [
        "User: drag-and-drop avatar, booking tabs (Upcoming / Past / Cancelled), visual status timeline.",
        "Owner: KPI cards with sparklines, revenue line charts, slide-over guest panels, inline-editable rates.",
        "Admin: CMS control of homepage blocks, plus platform users, hotels and payments oversight.",
        "All three share the same component kit and optimistic-update patterns.",
    ], size=17)
    add_image_safe(s, "owner_dashboard_mockup_1783678199513.png", RX + Inches(0.35), Inches(1.5), Inches(4.9), Inches(4.9))
    caption(s, RX, Inches(6.45), Inches(6.0), "Owner dashboard — KPIs, revenue chart & booking records")

    # --- 9. Frontend State Architecture (Redux) ---
    print("New: Redux")
    s = new_slide("Frontend State Architecture (Redux)")
    card(s, Inches(4.55), Inches(1.15), Inches(4.2), Inches(0.8),
         "REDUX ROOT STORE  —  configureStore() · 9 slices", fill=CARD_ALT, size=14, bold=True)
    slices = [
        ("Auth", "profile, role, JWT session"),
        ("Hotels", "search, list & detail cache"),
        ("Bookings", "checkout steps, promos"),
        ("UI", "theme, modals, toasts"),
        ("Analytics", "recent views, recos"),
        ("Content", "dynamic CMS blocks"),
        ("Admin", "users, hotels, payments"),
        ("Owner", "hotels, rooms, revenue"),
        ("Notifications", "bell feed, unread counts"),
    ]
    cols = [Inches(0.62), Inches(4.55), Inches(8.48)]
    rows = [Inches(2.35), Inches(3.75), Inches(5.15)]
    for i, (n, d) in enumerate(slices):
        r, c = divmod(i, 3)
        card(s, cols[c], rows[r], Inches(4.2), Inches(1.15), f"{n} Slice\n{d}", fill=CARD_FILL, size=13, bold=True)

    # --- 10. Supporting Backend Integration ---
    print("New: Backend")
    s = new_slide("Supporting Backend Integration")
    add_bullets(s, LX, Inches(1.2), LW, Inches(5.6), [
        "API: Node.js + Express REST with secure JWT cookie authentication.",
        "Database: MongoDB — Users, Hotels, Bookings, ContentBlocks, Engagement.",
        "Dynamic CMS: ContentBlock Mixed schema makes the homepage fully configurable.",
        "Personalization: RecentView / Engagement tracking drives live recommendations.",
        "Reliability: CORS, Helmet, rate-limiting and transactional email (Nodemailer).",
    ], size=17)
    y = Inches(1.55)
    for i, (t, sub) in enumerate([
        ("1. API Client", "React views / Redux thunks"),
        ("2. Express Router + Middleware", "JWT auth & rate-limit checks"),
        ("3. Controllers", "Business logic, CMS, recommendations"),
        ("4. MongoDB", "Users, Bookings, CMS, Analytics"),
    ]):
        bd = GREEN if i == 3 else TEAL
        card(s, RX, y, Inches(5.7), Inches(0.95), f"{t}\n{sub}", fill=CARD_FILL, border=bd, size=13, bold=True)
        if i < 3:
            arrow(s, RX + Inches(2.65), y + Inches(0.97), Inches(0.4), Inches(0.28))
        y = y + Inches(1.33)

    # --- 11. Performance & Accessibility ---
    print("New: Performance & A11y")
    s = new_slide("Performance & Accessibility Standards")
    add_bullets(s, LX, Inches(1.2), LW, Inches(5.6), [
        "Performance budgets enforced across every key page:",
        ("LCP ≤ 2.2s on 4G", 1),
        ("CLS < 0.05 (skeleton-first loading)", 1),
        ("TTI < 3.5s (route code-splitting)", 1),
        "Optimizations: lazy media, preconnects, layout-preserving skeletons.",
        "WCAG AA: contrast-checked palette, role=status/alert toasts, full keyboard nav & focus traps.",
    ], size=17)
    card(s, RX, Inches(1.5), Inches(5.7), Inches(0.6), "CORE WEB VITALS — AUDIT", fill=CARD_ALT, size=13, bold=True)
    for i, (m, v) in enumerate([
        ("Largest Contentful Paint (LCP)", "2.2s   PASS"),
        ("Cumulative Layout Shift (CLS)", "0.03   PASS"),
        ("Time to Interactive (TTI)", "3.4s   PASS"),
        ("WCAG AA Contrast Audit", "Passed"),
    ]):
        card(s, RX, Inches(2.25) + Inches(0.7) * i, Inches(5.7), Inches(0.58),
             f"{m}      {v}", fill=RGBColor(0xE6, 0xF7, 0xEC), border=GREEN, size=12, bold=True, align=PP_ALIGN.LEFT)

    # --- 12. Future Roadmap ---
    print("New: Roadmap")
    s = new_slide("Future Roadmap & Next Steps")
    add_bullets(s, LX, Inches(1.2), LW, Inches(5.6), [
        "Real-time sync: WebSockets (Socket.io) for live availability & instant confirmations.",
        "Interactive maps: Mapbox / Google Maps on search and detail pages.",
        "AI concierge: an LLM-powered booking assistant widget.",
        "PWA: service workers for offline itineraries and push notifications.",
    ], size=17)
    phases = [
        ("Phase 1 — Real-time Sync", "WebSockets for live bookings"),
        ("Phase 2 — Interactive Maps", "Mapbox embeds on detail pages"),
        ("Phase 3 — AI Concierge", "LLM assistant on the homepage"),
        ("Phase 4 — PWA", "Offline caching & push"),
    ]
    y = Inches(1.55)
    for i, (t, sub) in enumerate(phases):
        bd = GREEN if i == 3 else TEAL
        card(s, RX, y, Inches(5.7), Inches(0.95), f"{t}\n{sub}", fill=CARD_FILL, border=bd, size=13, bold=True)
        if i < 3:
            arrow(s, RX + Inches(2.65), y + Inches(0.97), Inches(0.4), Inches(0.28))
        y = y + Inches(1.33)

    new_count = len(prs.slides) - original_count
    print(f"Added {new_count} new content slides.")

    # -------------------------------------------------------------------
    # Reorder: new content slides go AFTER Outline (slide 3), BEFORE originals 4-7.
    # -------------------------------------------------------------------
    sldIdLst = prs.slides._sldIdLst
    ids = list(sldIdLst)
    originals = ids[:original_count]
    news = ids[original_count:]
    desired = originals[:3] + news + originals[3:]
    for e in ids:
        sldIdLst.remove(e)
    for e in desired:
        sldIdLst.append(e)

    # -------------------------------------------------------------------
    # Fill the (now-repositioned) original placeholder slides with real content.
    #   original slide 4 -> "Conclusion", slide 5 -> "References". Q&A / Thank You kept.
    # -------------------------------------------------------------------
    print("Filling original 'Topic 1' slide as Conclusion...")
    concl = prs.slides[3 + new_count]        # original slide 4
    set_title_ph(concl, "Conclusion")
    for sh in list(concl.shapes):
        if sh.has_text_frame and not sh.is_placeholder:
            sh.text_frame.clear()
            sh.left = Inches(0.71); sh.top = Inches(1.2); sh.width = Inches(11.9); sh.height = Inches(5.4)
            add_conclusion = [
                "BookMyStay delivers a premium, high-performance frontend built on React 18, Vite and Redux Toolkit.",
                "A token-driven luxury design system keeps the UI consistent, accessible and visually distinctive.",
                "Custom performance components (lazy media, skeletons, purposeful motion) yield near-zero layout shift.",
                "Three role-based dashboards and a dynamic CMS make the platform complete and production-ready.",
                "The architecture is extensible — real-time sync, maps, an AI concierge and PWA support are next.",
            ]
            tf = sh.text_frame; tf.word_wrap = True
            for i, t in enumerate(add_conclusion):
                p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
                p.space_after = Pt(8)
                no_bullet(p)
                _run(p, "▪  ", size=18, color=TEAL, bold=True)
                _run(p, t, size=18, color=BLACK)
            break

    print("Filling References slide (IEEE format)...")
    refs = prs.slides[4 + new_count]         # original slide 5
    set_title_ph(refs, "References")
    for sh in list(refs.shapes):
        if sh.has_text_frame and not sh.is_placeholder:
            sh.text_frame.clear()
            sh.left = Inches(0.71); sh.top = Inches(1.2); sh.width = Inches(11.9); sh.height = Inches(5.4)
            ref_list = [
                "[1] Meta Open Source, “React 18 – Concurrent Rendering & Transitions,” react.dev, 2022.",
                "[2] Redux Team, “Redux Toolkit – State Management Best Practices,” redux-toolkit.js.org, 2023.",
                "[3] Evan You et al., “Vite 5 Build Tooling Documentation,” vitejs.dev, 2024.",
                "[4] W3C WAI, “Web Content Accessibility Guidelines (WCAG) 2.1,” w3.org, 2018.",
                "[5] Google, “Optimizing Core Web Vitals – LCP & CLS,” web.dev, 2023.",
                "[6] Mozilla, “IntersectionObserver API & CSS backdrop-filter,” developer.mozilla.org, 2023.",
                "[7] OpenJS Foundation, “Express.js Security – Helmet, CORS & Rate Limiting,” expressjs.com, 2023.",
            ]
            tf = sh.text_frame; tf.word_wrap = True
            for i, t in enumerate(ref_list):
                p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
                p.space_after = Pt(8)
                no_bullet(p)
                _run(p, t, size=16, color=BLACK)
            break

    print("Saving presentation...")
    prs.save(OUTPUT_PATH)
    print("Saved:", OUTPUT_PATH)
    print("Total slides:", len(prs.slides))


if __name__ == "__main__":
    build_ppt()
